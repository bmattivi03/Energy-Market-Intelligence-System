#!/usr/bin/env python3
"""Assemble deck.pptx from the rendered deck.pdf: one full-bleed slide-image per page.

The deck is a bespoke HTML/CSS design, so the faithful PowerPoint is image-based — each
16:9 slide is the high-resolution render of the corresponding PDF page (pixel-identical to
deck.pdf; not editable text). Run build_pdf.sh first.

    python3 build_pptx.py        # -> presentation/deck.pptx
"""
import glob, os, pathlib, subprocess, sys, tempfile
from pptx import Presentation
from pptx.util import Inches

HERE = pathlib.Path(__file__).resolve().parent
PDF = HERE / "deck.pdf"
OUT = HERE / "deck.pptx"

if not PDF.exists():
    sys.exit("deck.pdf not found — run ./build_pdf.sh first")

with tempfile.TemporaryDirectory() as tmp:
    # 960x540pt page @ 144 dpi -> 1920x1080 px (exact 16:9)
    subprocess.run(["pdftoppm", "-png", "-r", "144", str(PDF), os.path.join(tmp, "s")],
                   check=True)
    pages = sorted(glob.glob(os.path.join(tmp, "s-*.png")))
    if not pages:
        sys.exit("no pages rendered")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    for png in pages:
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(png, 0, 0, width=prs.slide_width, height=prs.slide_height)
    prs.save(str(OUT))

print(f"wrote {OUT}  ({len(pages)} slides, {OUT.stat().st_size//1024} KB)")
